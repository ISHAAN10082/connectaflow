from __future__ import annotations

import json
import os
from io import BytesIO
from typing import Any

from loguru import logger


def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception:
        return ""
    reader = PdfReader(BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(parts)


def _extract_docx(data: bytes) -> str:
    try:
        import docx
    except Exception:
        return ""
    doc = docx.Document(BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return ""
    prs = Presentation(BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text
                if txt:
                    parts.append(txt)
    return "\n".join(parts)


def extract_text_from_file(filename: str, data: bytes) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return _extract_pdf(data)
    if ext == ".docx":
        return _extract_docx(data)
    if ext == ".pptx":
        return _extract_pptx(data)
    # Fallback: try utf-8 text
    try:
        return data.decode("utf-8", errors="ignore")
    except Exception:
        return ""


def compute_context_quality(fields: dict[str, Any]) -> int:
    required = [
        "company_name",
        "website_url",
        "product_description",
        "core_problem",
        "product_category",
        "pricing_model",
        "avg_deal_size",
        "customer_examples",
        "competitors",
        "geographic_focus",
    ]
    score = 0
    for key in required:
        val = fields.get(key)
        if isinstance(val, list):
            if len(val) > 0:
                score += 1
        elif isinstance(val, str):
            if val.strip():
                score += 1
        elif val is not None:
            score += 1
    return int(round((score / max(len(required), 1)) * 100))


async def parse_context_with_llm(text: str, extra: dict[str, Any], settings) -> dict[str, Any]:
    import litellm

    model = None
    api_key = None
    if settings.GROQ_API_KEY:
        model = "groq/llama-3.3-70b-versatile"
        api_key = settings.GROQ_API_KEY
    elif settings.GEMINI_API_KEY:
        model = "gemini/gemini-2.0-flash"
        api_key = settings.GEMINI_API_KEY
    else:
        raise ValueError("No LLM provider configured")

    prompt = f"""
You are extracting structured GTM context from notes, decks, and docs.
Return ONLY valid JSON with the following keys:
- company_name
- website_url
- product_description
- core_problem
- product_category
- pricing_model
- avg_deal_size
- customer_examples (array)
- competitors (array)
- geographic_focus (array)
- value_proposition
- sales_cycle_days
- decision_process
- key_integrations (array)
- why_customers_buy
- why_customers_churn
- common_objections (array)
- market_maturity
- context_notes (short summary)

If a field is unknown, use empty string or empty array.

EXISTING_HINTS (use to fill gaps, do not repeat verbatim):
{json.dumps(extra, default=str)}

DOCUMENT_TEXT:
{text[:20000]}
"""

    response = await litellm.acompletion(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        api_key=api_key,
        temperature=0.2,
        max_tokens=2000,
    )
    raw = response.choices[0].message.content.strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1]
        if raw.endswith("```"):
            raw = raw.rsplit("```", 1)[0]
    try:
        return json.loads(raw)
    except Exception:
        logger.error("Failed to parse context JSON")
        raise
